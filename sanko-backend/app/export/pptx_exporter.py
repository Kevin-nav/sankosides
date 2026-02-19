"""
PowerPoint Exporter

Generates PowerPoint files from RefinedSlide content with:
- Fully editable equations via OMML
- High-resolution diagram images
- Theme-aware styling
"""

import io
import base64
from typing import List, Optional, Tuple
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree
    HAS_PPTX = True
except ModuleNotFoundError:
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    MSO_ANCHOR = None
    MSO_SHAPE = None
    qn = None
    parse_xml = None
    etree = None
    HAS_PPTX = False

    def Inches(value):  # type: ignore
        return float(value)

    def Pt(value):  # type: ignore
        return float(value)

    def Emu(value):  # type: ignore
        return int(value)

from app.models.schemas import RefinedSlide, SlideContentType
from app.core.config import settings
try:
    from app.export.converters.latex_to_omml import latex_to_omml, OMML_NS
except ModuleNotFoundError:
    OMML_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"

    def latex_to_omml(_latex: str):
        return None, "LaTeX converter dependencies unavailable"
from app.core.logging import get_logger

logger = get_logger(__name__)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
SLIDE_WIDTH_EMU = 12192000   # 13.333" in EMU
SLIDE_HEIGHT_EMU = 6858000   # 7.5" in EMU

# Layout constants
MARGIN_LEFT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_BOTTOM = Inches(0.5)
TITLE_HEIGHT = Inches(1.0)
CONTENT_TOP = MARGIN_TOP + TITLE_HEIGHT + Inches(0.25)


def percent_to_emu(x_pct: float, y_pct: float, w_pct: float, h_pct: float):
    """Convert percentage-based coordinates into PPTX EMU values."""
    return (
        int(SLIDE_WIDTH_EMU * x_pct / 100),
        int(SLIDE_HEIGHT_EMU * y_pct / 100),
        int(SLIDE_WIDTH_EMU * w_pct / 100),
        int(SLIDE_HEIGHT_EMU * h_pct / 100),
    )


class PptxExporter:
    """
    Exports RefinedSlide content to PowerPoint format.
    
    Key features:
    - Editable equations via LaTeX → OMML conversion
    - Diagram embedding as high-res images
    - Theme-aware color and font settings
    """
    
    def __init__(
        self,
        theme_id: str = "academic",
        editable_equations: bool = True,
        image_dpi: int = 300,
    ):
        self.theme_id = theme_id
        self.editable_equations = editable_equations
        self.image_dpi = image_dpi
        
        # Theme color mapping (will be expanded with real theme system)
        self._theme_colors = self._get_theme_colors(theme_id) if HAS_PPTX else {}
    
    def _get_theme_colors(self, theme_id: str) -> dict:
        """Get color palette for theme."""
        if not HAS_PPTX:
            return {}

        themes = {
            "academic": {
                "primary": RGBColor(0x2D, 0x40, 0x59),      # Dark blue
                "secondary": RGBColor(0xEA, 0x54, 0x55),    # Accent red
                "background": RGBColor(0xFF, 0xFF, 0xFF),   # White
                "text": RGBColor(0x1A, 0x1A, 0x2E),         # Dark text
                "text_secondary": RGBColor(0x4A, 0x4A, 0x6A),
            },
            "modern": {
                "primary": RGBColor(0x66, 0x7E, 0xEA),      # Purple-blue
                "secondary": RGBColor(0xF0, 0x93, 0xFB),    # Pink accent
                "background": RGBColor(0xFF, 0xFF, 0xFF),
                "text": RGBColor(0x1F, 0x29, 0x37),
                "text_secondary": RGBColor(0x6B, 0x72, 0x80),
            },
            "dark": {
                "primary": RGBColor(0x00, 0xD4, 0xFF),      # Cyan
                "secondary": RGBColor(0xF0, 0x93, 0xFB),
                "background": RGBColor(0x0F, 0x0F, 0x23),   # Dark bg
                "text": RGBColor(0xCC, 0xCC, 0xCC),
                "text_secondary": RGBColor(0x88, 0x88, 0x88),
            },
        }
        return themes.get(theme_id, themes["academic"])
    
    def export(
        self,
        slides: List[RefinedSlide],
        title: str = "Presentation",
    ) -> bytes:
        """
        Export slides to PowerPoint format.
        
        Args:
            slides: List of RefinedSlide objects
            title: Presentation title
            
        Returns:
            PowerPoint file as bytes
        """
        if not HAS_PPTX:
            raise RuntimeError("python-pptx is required for PPTX export but is not installed")

        # Create presentation with 16:9 aspect ratio
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        
        # Add slides
        for slide in slides:
            self._add_slide(prs, slide)
        
        # Export to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        logger.info(f"Exported {len(slides)} slides to PowerPoint")
        return output.read()
    
    def _add_slide(self, prs: Presentation, slide: RefinedSlide):
        """Add a single slide to the presentation."""
        # Use blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        pptx_slide = prs.slides.add_slide(blank_layout)

        # New path: render from element tree only when export rollout is enabled.
        if settings.enable_element_tree_export and getattr(slide, "element_tree", None):
            self._render_element_tree_slide(pptx_slide, slide.element_tree)
            return
        
        # Route to appropriate renderer based on content type
        if slide.content_type == SlideContentType.TITLE:
            self._render_title_slide(pptx_slide, slide)
        elif slide.content_type == SlideContentType.EQUATION:
            self._render_equation_slide(pptx_slide, slide)
        elif slide.content_type == SlideContentType.DIAGRAM:
            self._render_diagram_slide(pptx_slide, slide)
        elif slide.content_type == SlideContentType.TWO_COLUMN:
            self._render_two_column_slide(pptx_slide, slide)
        elif slide.content_type == SlideContentType.CONCLUSION:
            self._render_conclusion_slide(pptx_slide, slide)
        else:
            # Default content slide
            self._render_content_slide(pptx_slide, slide)

    def _render_element_tree_slide(self, pptx_slide, tree):
        """Render a slide directly from structured element coordinates."""
        for element in getattr(tree, "elements", []) or []:
            left, top, width, height = percent_to_emu(
                element.x, element.y, element.width, element.height
            )
            element_type = element.type.value if hasattr(element.type, "value") else str(element.type)

            if element_type == "text":
                shape = pptx_slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
                frame = shape.text_frame
                frame.clear()
                paragraph = frame.paragraphs[0]
                runs = getattr(element.content, "runs", []) or []
                if not runs:
                    paragraph.text = ""
                    continue
                for idx, run_data in enumerate(runs):
                    run = paragraph.add_run()
                    run.text = str(getattr(run_data, "text", ""))
                    if getattr(run_data, "size", None):
                        run.font.size = Pt(int(run_data.size))
                    if getattr(run_data, "bold", False):
                        run.font.bold = True
                    if getattr(run_data, "italic", False):
                        run.font.italic = True
                continue

            if element_type == "image":
                image_url = getattr(element.content, "url", "") or ""
                # If we have an embedded data URL, decode and place it; otherwise placeholder.
                if image_url.startswith("data:image/") and ";base64," in image_url:
                    try:
                        b64 = image_url.split(";base64,", 1)[1]
                        raw = base64.b64decode(b64)
                        pptx_slide.shapes.add_picture(
                            io.BytesIO(raw),
                            Emu(left),
                            Emu(top),
                            width=Emu(width),
                            height=Emu(height),
                        )
                        continue
                    except Exception:
                        pass
                placeholder = pptx_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
                )
                placeholder.text_frame.text = "[Image]"
                continue

            if element_type in {"equation", "diagram"}:
                shape = pptx_slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
                frame = shape.text_frame
                frame.clear()
                paragraph = frame.paragraphs[0]
                paragraph.text = "[Equation]" if element_type == "equation" else "[Diagram]"
                paragraph.font.size = Pt(16)
                paragraph.font.italic = True
                continue
    
    def _render_title_slide(self, pptx_slide, slide: RefinedSlide):
        """Render a title slide."""
        # Center title
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(2.5),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = slide.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self._theme_colors["primary"]
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle (first bullet point if exists)
        if slide.bullet_points:
            subtitle_box = pptx_slide.shapes.add_textbox(
                MARGIN_LEFT,
                Inches(4.0),
                SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                Inches(1.0)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide.bullet_points[0]
            p.font.size = Pt(24)
            p.font.color.rgb = self._theme_colors["text_secondary"]
            p.alignment = PP_ALIGN.CENTER
    
    def _render_content_slide(self, pptx_slide, slide: RefinedSlide):
        """Render a standard content slide with title and bullets."""
        # Add title
        self._add_title(pptx_slide, slide.title)
        
        # Add bullet points
        if slide.bullet_points:
            self._add_bullets(
                pptx_slide,
                slide.bullet_points,
                left=MARGIN_LEFT,
                top=CONTENT_TOP,
                width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                height=SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM
            )
        
        # Add equation if present
        if slide.equation_latex:
            self._add_equation(pptx_slide, slide.equation_latex)
        
        # Add diagram if present
        if slide.diagram_svg:
            png_data = getattr(self, '_diagram_pngs', {}).get(slide.order)
            self._add_diagram(pptx_slide, slide.diagram_svg, png_data=png_data)
    
    def _render_equation_slide(self, pptx_slide, slide: RefinedSlide):
        """Render a slide focused on equations."""
        self._add_title(pptx_slide, slide.title)
        
        # Main equation in center
        if slide.equation_latex:
            self._add_equation(
                pptx_slide,
                slide.equation_latex,
                center=True,
                size="large"
            )
        
        # Supporting bullet points below
        if slide.bullet_points:
            self._add_bullets(
                pptx_slide,
                slide.bullet_points,
                left=MARGIN_LEFT,
                top=Inches(5.0),
                width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                height=Inches(2.0)
            )
    
    def _render_diagram_slide(self, pptx_slide, slide: RefinedSlide):
        """Render a slide focused on diagrams."""
        self._add_title(pptx_slide, slide.title)
        
        # Diagram in center
        if slide.diagram_svg:
            png_data = getattr(self, '_diagram_pngs', {}).get(slide.order)
            self._add_diagram(pptx_slide, slide.diagram_svg, center=True, png_data=png_data)
        
        # Caption/bullets below
        if slide.bullet_points:
            self._add_bullets(
                pptx_slide,
                slide.bullet_points,
                left=MARGIN_LEFT,
                top=Inches(5.5),
                width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                height=Inches(1.5)
            )
    
    def _render_two_column_slide(self, pptx_slide, slide: RefinedSlide):
        """Render a two-column slide."""
        self._add_title(pptx_slide, slide.title)
        
        # Split bullets into two columns
        mid = len(slide.bullet_points) // 2
        left_bullets = slide.bullet_points[:mid]
        right_bullets = slide.bullet_points[mid:]
        
        col_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.5)) / 2
        
        # Left column
        if left_bullets:
            self._add_bullets(
                pptx_slide,
                left_bullets,
                left=MARGIN_LEFT,
                top=CONTENT_TOP,
                width=col_width,
                height=SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM
            )
        
        # Right column
        if right_bullets:
            self._add_bullets(
                pptx_slide,
                right_bullets,
                left=MARGIN_LEFT + col_width + Inches(0.5),
                top=CONTENT_TOP,
                width=col_width,
                height=SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM
            )
    
    def _render_conclusion_slide(self, pptx_slide, slide: RefinedSlide):
        """Render a conclusion/summary slide."""
        # Add "Key Takeaways" header
        self._add_title(pptx_slide, slide.title or "Key Takeaways")
        
        # Larger bullets for emphasis
        if slide.bullet_points:
            self._add_bullets(
                pptx_slide,
                slide.bullet_points,
                left=MARGIN_LEFT,
                top=CONTENT_TOP,
                width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                height=SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM,
                font_size=Pt(22)
            )
    
    def _add_title(self, pptx_slide, title: str):
        """Add title to slide."""
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            MARGIN_TOP,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            TITLE_HEIGHT
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._theme_colors["primary"]
    
    def _add_bullets(
        self,
        pptx_slide,
        bullets: List[str],
        left,
        top,
        width,
        height,
        font_size=Pt(18)
    ):
        """Add bullet points to slide."""
        text_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {bullet}"
            p.font.size = font_size
            p.font.color.rgb = self._theme_colors["text"]
            p.space_after = Pt(8)
    
    def _add_equation(
        self,
        pptx_slide,
        latex: str,
        center: bool = False,
        size: str = "normal"
    ):
        """
        Add equation to slide.
        
        If editable_equations is True, attempts OMML conversion.
        Falls back to placeholder text if conversion fails.
        """
        if self.editable_equations:
            omml, error = latex_to_omml(latex)
            
            if omml is not None:
                # Successfully converted - add as OMML
                self._insert_omml_equation(pptx_slide, omml, center, size)
                return
            else:
                logger.warning(f"OMML conversion failed: {error}, using placeholder")
        
        # Fallback: add as text placeholder
        # In production, this would embed SVG as image
        eq_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT if not center else Inches(2),
            Inches(3) if center else CONTENT_TOP,
            Inches(9) if center else Inches(6),
            Inches(1.5)
        )
        tf = eq_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Equation: {latex[:50]}...]" if len(latex) > 50 else f"[Equation: {latex}]"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self._theme_colors["text_secondary"]
        if center:
            p.alignment = PP_ALIGN.CENTER
    
    def _insert_omml_equation(self, pptx_slide, omml, center: bool, size: str):
        """Insert OMML equation into slide."""
        # Create a text box to hold the equation
        left = Inches(2) if center else MARGIN_LEFT
        top = Inches(3) if center else Inches(4)
        width = Inches(9) if center else Inches(6)
        height = Inches(1.5)
        
        text_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        
        if center:
            p.alignment = PP_ALIGN.CENTER
        
        # Get the paragraph XML element
        p_elem = p._p
        
        # Wrap OMML in oMathPara for proper display
        omml_para = etree.Element(f"{{{OMML_NS}}}oMathPara")
        omml_para.append(omml)
        
        # Insert into paragraph
        # Note: This is a simplified approach. Full integration requires
        # more complex XML manipulation with the Word/PowerPoint namespaces.
        try:
            p_elem.append(omml_para)
            logger.debug("Inserted OMML equation into slide")
        except Exception as e:
            logger.warning(f"Failed to insert OMML: {e}")
            # Fallback to text
            p.text = "[Equation - see original LaTeX]"
    
    def _add_diagram(
        self,
        pptx_slide,
        svg_content: str,
        center: bool = False,
        png_data: Optional[bytes] = None,
    ):
        """
        Add diagram to slide.
        
        If png_data is provided, embeds as high-res image.
        Otherwise adds a placeholder shape.
        
        Args:
            pptx_slide: PowerPoint slide object
            svg_content: SVG content (for reference/fallback)
            center: Whether to center the diagram
            png_data: Pre-rendered PNG bytes (optional)
        """
        left = Inches(3) if center else MARGIN_LEFT
        top = Inches(2.5) if center else CONTENT_TOP
        width = Inches(7) if center else Inches(5)
        height = Inches(3)
        
        if png_data:
            # Embed PNG image
            try:
                image_stream = io.BytesIO(png_data)
                pptx_slide.shapes.add_picture(
                    image_stream,
                    left, top,
                    width=width,
                    # Height is auto-calculated to maintain aspect ratio
                )
                logger.debug(f"Embedded diagram PNG: {len(png_data)} bytes")
                return
            except Exception as e:
                logger.warning(f"Failed to embed PNG: {e}, using placeholder")
        
        # Fallback: Add placeholder shape
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)  # Light gray
        
        # Add text indicating diagram
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[Diagram - render service unavailable]"
        p.font.size = Pt(14)
        p.font.color.rgb = self._theme_colors["text_secondary"]
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE


def export_to_pptx(
    slides: List[RefinedSlide],
    title: str = "Presentation",
    theme_id: str = "academic",
    editable_equations: bool = True,
    image_dpi: int = 300,
    diagram_pngs: Optional[dict] = None,
) -> bytes:
    """
    Export slides to PowerPoint format.
    
    Args:
        slides: List of RefinedSlide objects
        title: Presentation title
        theme_id: Theme to apply
        editable_equations: Whether to use OMML for editable equations
        image_dpi: DPI for embedded images
        diagram_pngs: Optional dict mapping slide order -> PNG bytes for diagrams
        
    Returns:
        PowerPoint file as bytes
    """
    exporter = PptxExporter(
        theme_id=theme_id,
        editable_equations=editable_equations,
        image_dpi=image_dpi,
    )
    
    # If we have pre-rendered diagrams, store them on the exporter
    if diagram_pngs:
        exporter._diagram_pngs = diagram_pngs
    else:
        exporter._diagram_pngs = {}
    
    return exporter.export(slides, title)


async def export_to_pptx_async(
    slides: List[RefinedSlide],
    title: str = "Presentation",
    theme_id: str = "academic",
    editable_equations: bool = True,
    image_dpi: int = 300,
) -> bytes:
    """
    Export slides to PowerPoint format with async diagram rendering.
    
    This function pre-renders any SVG diagrams to PNG using the render
    service before generating the PowerPoint file.
    
    Args:
        slides: List of RefinedSlide objects
        title: Presentation title
        theme_id: Theme to apply
        editable_equations: Whether to use OMML for editable equations
        image_dpi: DPI for embedded images
        
    Returns:
        PowerPoint file as bytes
    """
    from app.export.render_client import get_render_client
    
    render_client = get_render_client()
    
    # Pre-render diagrams
    diagram_pngs = {}
    
    for slide in slides:
        # Handle SVG diagrams
        if slide.diagram_svg:
            logger.info(f"Pre-rendering SVG diagram for slide {slide.order}")
            png_data, error = await render_client.svg_to_png(
                slide.diagram_svg,
                width=800,
                height=600,
                scale=2,  # 2x for high-res
            )
            
            if png_data:
                diagram_pngs[slide.order] = png_data
                logger.debug(f"Diagram rendered: {len(png_data)} bytes")
            else:
                logger.warning(f"Failed to render SVG diagram for slide {slide.order}: {error}")
        
        # Handle Mermaid source (needs two-step: Mermaid -> SVG -> PNG)
        elif hasattr(slide, 'diagram_mermaid') and slide.diagram_mermaid:
            logger.info(f"Pre-rendering Mermaid diagram for slide {slide.order}")
            png_data, error = await render_client.render_mermaid_to_png(
                slide.diagram_mermaid,
                width=800,
                height=500,
                scale=2,
            )
            
            if png_data:
                diagram_pngs[slide.order] = png_data
                logger.debug(f"Mermaid diagram rendered: {len(png_data)} bytes")
            else:
                logger.warning(f"Failed to render Mermaid diagram for slide {slide.order}: {error}")
    
    # Generate PPTX with pre-rendered diagrams
    return export_to_pptx(
        slides=slides,
        title=title,
        theme_id=theme_id,
        editable_equations=editable_equations,
        image_dpi=image_dpi,
        diagram_pngs=diagram_pngs,
    )
